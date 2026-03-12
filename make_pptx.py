import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation(language):
    prs = Presentation()
    
    # Define colors
    BG_COLOR = RGBColor(30, 30, 46)  # Dark Blue/Slate
    TEXT_COLOR = RGBColor(255, 255, 255)
    ACCENT_COLOR = RGBColor(0, 255, 255)  # Cyan
    
    content = {
        'it': {
            'filename': 'Parvis Loggy-Eye_Presentation_IT.pptx',
            'slides': [
                {
                    'title': 'Parvis Loggy-Eye: Diagnostic Intelligence',
                    'subtitle': 'L\'integrazione forense per l\'Ecosistema Abaco Server',
                    'type': 'title'
                },
                {
                    'title': 'Sfide e Complessità Attuali',
                    'points': [
                        'Dati frammentati tra ICS e Server Library.',
                        'File di grandi dimensioni (>100MB) complessi da navigare.',
                        'MTTR elevato dovuto a identificazione lenta del root cause.'
                    ]
                },
                {
                    'title': 'La Soluzione: Performance & UI',
                    'points': [
                        'Architettura Standalone basata su Electron/React.',
                        'Rendering virtualizzato ad alta fluidità.',
                        'Parsing asincrono per un\'esperienza d\'uso premium.'
                    ]
                },
                {
                    'title': 'Motore di Analisi Avanzato',
                    'points': [
                        'Boot Tracking: Rilevamento automatico di anomalie e blackout.',
                        'Dashboard Integrazione: Ogni evento è cliccabile per il jump istantaneo al log.',
                        'MDR Ready: Estrazione automatica locale e remota.'
                    ]
                },
                {
                    'title': 'Valore Aziendale',
                    'points': [
                        'Field Support: Diagnostica immediata senza dipendenza da sviluppatori.',
                        'Report Professionali: Comunicazione chiara verso il cliente finale.',
                        'QA Optimization: Validazione flussi di produzione accelerata.'
                    ]
                },
                {
                    'title': 'Conclusioni',
                    'points': [
                        'Standardizzazione della diagnostica Abaco.',
                        'Repo GitHub attivo e tool pronto al deployment.',
                        'Migliore percezione tecnologica del brand.'
                    ]
                }
            ],
            'footer': '© 2026 Parvis Loggy-Eye Analysis Suite — Progettato per l\'Eccellenza.'
        },
        'en': {
            'filename': 'Parvis Loggy-Eye_Presentation_EN.pptx',
            'slides': [
                {
                    'title': 'Parvis Loggy-Eye: Diagnostic Intelligence',
                    'subtitle': 'Forensic integration for the Abaco Server Ecosystem',
                    'type': 'title'
                },
                {
                    'title': 'Current Challenges & Complexity',
                    'points': [
                        'Fragmented data between ICS and Server Library.',
                        'Large files (>100MB) complex to navigate.',
                        'High MTTR due to slow root cause identification.'
                    ]
                },
                {
                    'title': 'The Solution: Performance & UI',
                    'points': [
                        'Standalone architecture based on Electron/React.',
                        'High-fluidity virtualized rendering.',
                        'Asynchronous parsing for a premium user experience.'
                    ]
                },
                {
                    'title': 'Advanced Analysis Engine',
                    'points': [
                        'Boot Tracking: Automatic detection of anomalies and blackouts.',
                        'Integration Dashboard: Every event is clickable for instant jump to log.',
                        'MDR Ready: Automatic local and remote extraction.'
                    ]
                },
                {
                    'title': 'Business Value',
                    'points': [
                        'Field Support: Immediate diagnostics without developer dependency.',
                        'Professional Reports: Clear communication with the end client.',
                        'QA Optimization: Accelerated production flow validation.'
                    ]
                },
                {
                    'title': 'Conclusion',
                    'points': [
                        'Standardization of Abaco diagnostics.',
                        'Active GitHub repo and tool ready for deployment.',
                        'Improved technological perception of the brand.'
                    ]
                }
            ],
            'footer': '© 2026 Parvis Loggy-Eye Analysis Suite — Designed for Excellence.'
        }
    }
    
    lang_content = content[language]
    
    for slide_data in lang_content['slides']:
        if slide_data.get('type') == 'title':
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = slide_data['title']
            subtitle.text = slide_data['subtitle']
            
            # Styling
            for p in title.text_frame.paragraphs:
                p.font.color.rgb = ACCENT_COLOR
                p.font.bold = True
                p.font.size = Pt(44)
            
            for p in subtitle.text_frame.paragraphs:
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(24)
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_COLOR
            
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_data['title']
            
            # Title styling
            for p in title.text_frame.paragraphs:
                p.font.color.rgb = ACCENT_COLOR
                p.font.bold = True
                p.font.size = Pt(36)
            
            # Body points
            tf = body.text_frame
            tf.text = slide_data['points'][0]
            for point in slide_data['points'][1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
            for p in tf.paragraphs:
                p.font.color.rgb = TEXT_COLOR
                p.font.size = Pt(20)
                p.space_after = Pt(12)

        # Add Footer
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = lang_content['footer']
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(120, 120, 120)
        p.alignment = PP_ALIGN.CENTER

    prs.save(lang_content['filename'])
    print(f"Generated {lang_content['filename']}")

if __name__ == "__main__":
    create_presentation('it')
    create_presentation('en')
